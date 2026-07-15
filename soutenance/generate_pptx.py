# -*- coding: utf-8 -*-
"""Génère la présentation PowerPoint de soutenance de Market BF (10 min, binôme)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = os.path.join(os.path.dirname(__file__), "Soutenance_MarketBF.pptx")
AUTHORS = "NITIEMA Martial  •  [Nom du binôme]"

GREEN_DARK = RGBColor(0x0E, 0x3D, 0x20)
GREEN = RGBColor(0x16, 0x65, 0x34)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
RED = RGBColor(0xC0, 0x39, 0x2B)
YELLOW = RGBColor(0xF0, 0xB4, 0x29)
GREY = RGBColor(0x37, 0x41, 0x51)
GREY_LIGHT = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xBB, 0xF7, 0xD0)
BG = RGBColor(0xFA, 0xFD, 0xFB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(s, x, y, w, h, text, size=18, color=GREY, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    r.font.name = font
    return tb

def bullets(s, x, y, w, h, items, size=16, color=GREY, gap=6, marker="•  "):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            head, rest = it
        else:
            head, rest = None, it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = marker; r.font.size = Pt(size)
        r.font.color.rgb = GREEN; r.font.bold = True
        if head:
            r2 = p.add_run(); r2.text = head + " — "
            r2.font.size = Pt(size); r2.font.bold = True; r2.font.color.rgb = GREEN_DARK
        r3 = p.add_run(); r3.text = rest
        r3.font.size = Pt(size); r3.font.color.rgb = color
    return tb

def header(s, kicker, title):
    rect(s, 0, 0, SW, Inches(0.14), GREEN)
    rect(s, 0, Inches(0.14), Inches(2.2), Inches(0.05), YELLOW)
    txt(s, Inches(0.55), Inches(0.34), SW - Inches(1.1), Inches(0.4),
        kicker.upper(), size=12, color=RED, bold=True)
    txt(s, Inches(0.55), Inches(0.68), SW - Inches(1.1), Inches(0.9),
        title, size=30, color=GREEN_DARK, bold=True)

def footer(s, n):
    txt(s, Inches(0.55), SH - Inches(0.48), Inches(6), Inches(0.35),
        "Market BF — Soutenance", size=10, color=GREY_LIGHT)
    txt(s, SW - Inches(1.3), SH - Inches(0.48), Inches(0.8), Inches(0.35),
        str(n), size=10, color=GREY_LIGHT, align=PP_ALIGN.RIGHT)

def card(s, x, y, w, h, title, lines, title_color=GREEN_DARK, fill=WHITE, tsize=15, bsize=12.5):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = RGBColor(0xD6, 0xE5, 0xDB); sh.line.width = Pt(1)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(tsize); r.font.bold = True; r.font.color.rgb = title_color
    p.space_after = Pt(5)
    for ln in lines:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = ln
        r2.font.size = Pt(bsize); r2.font.color.rgb = GREY
        p2.space_after = Pt(3)
    return sh

def table(s, x, y, w, headers, rows, col_ratios=None, size=12, header_size=12.5, row_h=0.42):
    n_rows, n_cols = len(rows) + 1, len(headers)
    shp = s.shapes.add_table(n_rows, n_cols, x, y, w, Inches(row_h * n_rows))
    t = shp.table
    if col_ratios:
        total = sum(col_ratios)
        for i, ratio in enumerate(col_ratios):
            t.columns[i].width = Emu(int(w * ratio / total))
    for j, htxt in enumerate(headers):
        c = t.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = GREEN
        c.margin_top = c.margin_bottom = Inches(0.03)
        p = c.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htxt
        r.font.bold = True; r.font.size = Pt(header_size); r.font.color.rgb = WHITE
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else GREEN_LIGHT
            c.margin_top = c.margin_bottom = Inches(0.03)
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.size = Pt(size); r.font.color.rgb = GREY
    return t

n = 0
# ------------------------------------------------------------------ 1. titre
s = slide()
rect(s, 0, 0, SW, SH, GREEN_DARK)
rect(s, 0, 0, SW, Inches(0.35), RED)
rect(s, 0, Inches(0.35), SW, Inches(0.09), YELLOW)
rect(s, 0, SH - Inches(0.35), SW, Inches(0.35), GREEN)
txt(s, Inches(1), Inches(1.85), SW - Inches(2), Inches(1.2), "Market BF",
    size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.05), SW - Inches(2), Inches(0.6),
    "Marketplace locale et gestion commerciale pour les commerces du Burkina Faso",
    size=21, color=PALE, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.15), SW - Inches(2), Inches(0.5),
    "SOUTENANCE DE PROJET — 10 MINUTES", size=17, color=YELLOW, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.1), SW - Inches(2), Inches(0.45),
    "Application web full-stack  •  Node.js / Express  •  React  •  SQLite (Turso)",
    size=14, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.55), SW - Inches(2), Inches(0.45),
    "En production : https://marketbf.onrender.com", size=14, color=WHITE,
    align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.35), SW - Inches(2), Inches(0.4),
    "Présenté par %s — Juillet 2026" % AUTHORS, size=14, color=PALE,
    align=PP_ALIGN.CENTER)

# ------------------------------------------------------------------- 2. plan
n += 1
s = slide(); header(s, "Sommaire", "Plan — 10 minutes à deux voix"); footer(s, n)
card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(4.3),
     "Partie 1 — Intervenant 1   (≈ 5 min)",
     ["1.  Contexte et problématique",
      "2.  Market BF en bref",
      "3.  Les trois espaces : client, commerçant, admin",
      "4.  Le parcours d'une commande",
      "5.  Paiements adaptés au contexte local"], tsize=17, bsize=15)
card(s, Inches(6.8), Inches(1.95), Inches(6.0), Inches(4.3),
     "Partie 2 — Intervenant 2   (≈ 5 min)",
     ["6.  Stack technique et architecture",
      "7.  Sécurité",
      "8.  Déploiement en production",
      "9.  Démonstration en direct",
      "10.  Conclusion et perspectives"], title_color=GREEN, tsize=17, bsize=15)

# --------------------------------------------------------------- 3. contexte
n += 1
s = slide(); header(s, "Partie 1 • Introduction", "Contexte et problématique"); footer(s, n)
card(s, Inches(0.55), Inches(1.85), Inches(6.0), Inches(2.35), "Le constat",
     ["Au Burkina Faso, les commerces de proximité gèrent stocks, ventes et "
      "commandes sans outil numérique.",
      "Les clients se déplacent de boutique en boutique, sans visibilité sur "
      "la disponibilité des produits ni sur les prix."], tsize=17, bsize=14)
card(s, Inches(6.8), Inches(1.85), Inches(6.0), Inches(2.35), "La réponse : Market BF",
     ["Une plateforme unique qui digitalise la gestion des commerçants "
      "(stock, commandes, statistiques)…",
      "…et offre aux clients une marketplace géolocalisée avec les moyens "
      "de paiement locaux."], title_color=GREEN, tsize=17, bsize=14)
rect(s, Inches(0.55), Inches(4.6), SW - Inches(1.1), Inches(1.9), GREEN_DARK)
txt(s, Inches(0.9), Inches(4.85), SW - Inches(1.8), Inches(0.5),
    "Le projet en chiffres", size=16, color=YELLOW, bold=True)
stats = [("≈ 7 000", "lignes de code"), ("9", "modules d'API REST"), ("16", "pages d'interface"),
         ("3", "rôles utilisateurs"), ("9", "tables de données"), ("En ligne", "sur Render")]
x = Inches(0.9)
for val, lab in stats:
    txt(s, x, Inches(5.3), Inches(1.95), Inches(0.5), val, size=24, color=WHITE, bold=True)
    txt(s, x, Inches(5.85), Inches(1.95), Inches(0.4), lab, size=11.5, color=PALE)
    x += Inches(1.98)

# ----------------------------------------------------------------- 4. acteurs
n += 1
s = slide(); header(s, "Partie 1 • Fonctionnel", "Trois acteurs, trois espaces dédiés"); footer(s, n)
card(s, Inches(0.55), Inches(1.95), Inches(4.0), Inches(4.6), "Client",
     ["Catalogue et recherche par nom / catégorie", "Carte interactive des boutiques (Leaflet)",
      "Disponibilité des produits en temps réel", "Panier et commande en ligne",
      "Paiement : Orange Money, Moov, Wave, à la livraison", "Reçu et confirmation par e-mail",
      "Historique et suivi des commandes", "Avis et notation des boutiques (1 à 5)"],
     tsize=18, bsize=13)
card(s, Inches(4.7), Inches(1.95), Inches(4.0), Inches(4.6), "Commerçant",
     ["Création et profil de boutique géolocalisée", "Gestion des produits avec photos",
      "Gestion du stock : quantités et seuils", "Alertes de stock faible",
      "Mouvements d'entrée / sortie tracés", "Traitement des commandes par statut",
      "Statistiques : CA, volumes, top produits", "Graphiques interactifs (Recharts)"],
     title_color=GREEN, tsize=18, bsize=13)
card(s, Inches(8.85), Inches(1.95), Inches(4.0), Inches(4.6), "Administrateur",
     ["Validation ou rejet des boutiques", "Workflow pending / active / rejected",
      "Gestion des utilisateurs", "Statistiques globales de la plateforme",
      "Garant de la qualité et de la sécurité"],
     title_color=RED, tsize=18, bsize=13)

# ---------------------------------------------------------------- 5. parcours
n += 1
s = slide(); header(s, "Partie 1 • Fonctionnel", "Le parcours d'une commande"); footer(s, n)
steps = [
    ("1. Commande", "Le client remplit son panier et choisit son mode de paiement — il reçoit un reçu par e-mail"),
    ("2. Stock", "Le stock est décrémenté automatiquement, de façon transactionnelle : pas de survente possible"),
    ("3. Traitement", "Le commerçant voit la commande arriver et la confirme depuis son tableau de bord"),
    ("4. Livraison", "Statut « livrée » : paiement complété, e-mail de confirmation envoyé au client"),
    ("5. Avis", "Le client peut noter la boutique et laisser un commentaire"),
]
y = Inches(1.95)
for t_, d_ in steps:
    card(s, Inches(0.55), y, Inches(12.25), Inches(0.82), "", [], fill=WHITE)
    txt(s, Inches(0.8), y + Inches(0.18), Inches(2.2), Inches(0.5), t_,
        size=16, color=GREEN_DARK, bold=True)
    txt(s, Inches(3.1), y + Inches(0.2), Inches(9.5), Inches(0.5), d_, size=14, color=GREY)
    y += Inches(0.92)
txt(s, Inches(0.55), y + Inches(0.05), Inches(12.25), Inches(0.5),
    "En cas d'annulation, le stock est automatiquement restitué.", size=13,
    color=GREY_LIGHT)

# ---------------------------------------------------------------- 6. paiement
n += 1
s = slide(); header(s, "Partie 1 • Fonctionnel", "Des paiements adaptés au contexte local"); footer(s, n)
card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(2.4), "4 méthodes proposées",
     ["Orange Money", "Moov Money", "Wave", "Paiement à la livraison"],
     title_color=GREEN, bsize=14.5)
card(s, Inches(6.8), Inches(1.95), Inches(6.0), Inches(2.4), "Suivi complet des transactions",
     ["Chaque commande génère un enregistrement de paiement",
      "Méthode, montant et statut conservés",
      "Statuts contrôlés : pending → completed / failed",
      "Reçu détaillé envoyé par e-mail"], bsize=13.5)
rect(s, Inches(0.55), Inches(4.75), SW - Inches(1.1), Inches(1.85), GREEN_LIGHT)
txt(s, Inches(0.9), Inches(4.9), SW - Inches(1.8), Inches(0.45),
    "Pourquoi des paiements mobiles simulés ?", size=15, color=GREEN_DARK, bold=True)
txt(s, Inches(0.9), Inches(5.35), SW - Inches(1.8), Inches(1.2),
    "Les API Orange Money / Moov / Wave exigent un contrat marchand et un agrégateur agréé. "
    "Le flux complet (choix de la méthode, enregistrement, statuts, validation à la livraison) "
    "est implémenté : brancher un agrégateur réel ne demandera aucune refonte.",
    size=13.5, color=GREY)

# ------------------------------------------------------------ 7. architecture
n += 1
s = slide(); header(s, "Partie 2 • Technique", "Stack et architecture"); footer(s, n)
card(s, Inches(0.55), Inches(1.95), Inches(3.7), Inches(2.6), "Frontend — SPA React",
     ["React 18 + Vite + Tailwind CSS", "16 pages / 3 espaces",
      "Leaflet (carte), Recharts (graphiques)", "Routage protégé par rôle"], bsize=12.5)
card(s, Inches(4.85), Inches(1.95), Inches(3.7), Inches(2.6), "API REST — Express",
     ["Node.js + Express", "9 modules de routes (/api/…)",
      "Middleware JWT + contrôle des rôles", "Validation systématique des entrées"],
     title_color=GREEN, bsize=12.5)
card(s, Inches(9.15), Inches(1.95), Inches(3.7), Inches(2.6), "Données & services",
     ["SQLite local / Turso en production", "Cloudinary : images persistantes",
      "SMTP : e-mails transactionnels", "Google Identity : connexion Google"],
     title_color=RED, bsize=12.5)
txt(s, Inches(4.3), Inches(2.9), Inches(0.6), Inches(0.5), "⇄", size=28, color=GREEN_DARK, bold=True)
txt(s, Inches(8.6), Inches(2.9), Inches(0.6), Inches(0.5), "⇄", size=28, color=GREEN_DARK, bold=True)
rect(s, Inches(0.55), Inches(4.95), SW - Inches(1.1), Inches(1.7), GREEN_LIGHT)
txt(s, Inches(0.9), Inches(5.1), SW - Inches(1.8), Inches(0.4),
    "Fil conducteur des choix techniques", size=14, color=GREEN_DARK, bold=True)
txt(s, Inches(0.9), Inches(5.55), SW - Inches(1.8), Inches(1.0),
    "Une plateforme utilisable partout au Burkina Faso : web responsive (aucune installation), "
    "cartographie OpenStreetMap sans clé API, hébergement et base de données gratuits — "
    "coût d'infrastructure : 0 F CFA.", size=13.5, color=GREY)

# ---------------------------------------------------------------- 8. sécurité
n += 1
s = slide(); header(s, "Partie 2 • Technique", "Sécurité : un durcissement complet"); footer(s, n)
card(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(2.15), "Authentification",
     ["JWT HS256, secret fort obligatoire en production", "bcrypt 12 rounds + politique de mot de passe",
      "Anti-énumération de comptes (temps constant)", "Réinitialisation : token 256 bits haché, 1 h, usage unique"],
     bsize=12.5)
card(s, Inches(6.8), Inches(1.95), Inches(6.0), Inches(2.15), "Protection de l'API",
     ["Helmet : en-têtes HTTP de sécurité", "Rate limiting : 300 req/15 min, 20 sur l'auth",
      "CORS restreint, payload JSON limité à 100 ko", "express-validator sur toutes les routes"],
     title_color=GREEN, bsize=12.5)
card(s, Inches(0.55), Inches(4.35), Inches(6.0), Inches(2.15), "Données & uploads",
     ["Transactions atomiques sur le stock", "Uploads : extension + type MIME vérifiés, 5 Mo max",
      "Noms de fichiers aléatoires", "Erreurs internes jamais exposées au client"],
     title_color=RED, bsize=12.5)
card(s, Inches(6.8), Inches(4.35), Inches(6.0), Inches(2.15), "Persistance des données",
     ["Base répliquée et persistante sur Turso", "Images persistantes sur Cloudinary",
      "Secrets en variables d'environnement", "Un crash ou redéploiement ne perd aucune donnée"],
     title_color=GREEN_DARK, bsize=12.5)

# -------------------------------------------------------------- 9. déploiement
n += 1
s = slide(); header(s, "Partie 2 • Production", "Déploiement continu, coût : 0 F CFA"); footer(s, n)
steps = [("1. GitHub", "git push sur la branche main"),
         ("2. Render", "build Docker automatique (render.yaml)"),
         ("3. En ligne", "https://marketbf.onrender.com")]
x = Inches(0.55)
for t_, d_ in steps:
    card(s, x, Inches(2.0), Inches(4.0), Inches(1.5), t_, [d_], bsize=13)
    x += Inches(4.15)
bullets(s, Inches(0.9), Inches(4.0), Inches(11.8), Inches(2.6), [
    ("Docker", "build multi-étapes : build Vite du frontend puis image Node unique servant API + site (pas de CORS)"),
    ("Turso", "base SQLite hébergée — comptes, produits et commandes survivent aux redémarrages du plan gratuit"),
    ("Cloudinary", "stockage persistant des images uploadées (photos produits, logos)"),
    ("Secrets", "JWT, SMTP, OAuth Google en variables d'environnement Render — jamais dans le code"),
    ("Local", "./start.sh ou docker-compose up : SQLite et dossier uploads locaux, aucun service externe requis"),
], size=14.5, gap=10)

# --------------------------------------------------------------------- 10. démo
n += 1
s = slide(); header(s, "Partie 2 • Démonstration", "Démonstration en direct"); footer(s, n)
bullets(s, Inches(0.9), Inches(2.0), Inches(11.8), Inches(2.9), [
    ("Client", "recherche un produit, consulte la carte, commande avec Orange Money → reçu par e-mail"),
    ("Commerçant", "voit la commande arriver, la confirme puis la marque livrée → stock décrémenté, alerte si seuil atteint"),
    ("Administrateur", "valide une boutique en attente → elle devient visible dans le catalogue et sur la carte"),
], size=17, gap=14)
rect(s, Inches(0.55), Inches(4.55), SW - Inches(1.1), Inches(2.0), GREEN_DARK)
txt(s, Inches(0.9), Inches(4.75), Inches(11.8), Inches(0.45),
    "https://marketbf.onrender.com", size=20, color=YELLOW, bold=True)
txt(s, Inches(0.9), Inches(5.3), Inches(11.8), Inches(1.1),
    "Comptes de démonstration :  admin@marketbf.com / Admin123!   •   "
    "merchant1@marketbf.com / Merchant1!   •   client@marketbf.com / Client123!",
    size=14, color=WHITE)
txt(s, Inches(0.9), Inches(5.95), Inches(11.8), Inches(0.5),
    "Astuce plan gratuit : ouvrir le site 2 minutes avant la soutenance (réveil du service ~1 min).",
    size=12, color=PALE)

# --------------------------------------------------------------- 11. conclusion
n += 1
s = slide(); header(s, "Conclusion", "Bilan et perspectives"); footer(s, n)
card(s, Inches(0.55), Inches(1.9), Inches(6.0), Inches(3.4), "Ce qui est livré",
     ["Une plateforme complète et en production", "3 espaces fonctionnels : client, commerçant, admin",
      "Chaîne complète : recherche → commande → livraison", "Données persistantes, sécurité durcie",
      "Coût d'infrastructure : 0 F CFA"], bsize=14)
card(s, Inches(6.8), Inches(1.9), Inches(6.0), Inches(3.4), "Les prochaines étapes",
     ["Paiements mobiles réels via un agrégateur", "Factures PDF téléchargeables",
      "Module de livraison intégré", "Application mobile (PWA / React Native)",
      "Recommandations par IA, fidélité, wallet"], title_color=GREEN, bsize=14)
rect(s, 0, Inches(5.6), SW, Inches(1.9), GREEN_DARK)
txt(s, Inches(1), Inches(6.0), SW - Inches(2), Inches(0.6),
    "Merci de votre attention — place à vos questions",
    size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.65), SW - Inches(2), Inches(0.4),
    "Market BF  •  https://marketbf.onrender.com", size=14, color=YELLOW,
    align=PP_ALIGN.CENTER)

prs.save(OUT)
print("PPTX généré :", OUT, "-", len(prs.slides._sldIdLst), "diapositives")
